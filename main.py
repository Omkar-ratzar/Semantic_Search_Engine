from sentence_transformers import SentenceTransformer
import numpy as np
import json
import docx2txt
from db_connection import mark_processing,mark_processed,get_file_id_by_path,get_image_text_for_embedding
# from image_desc import extract_image
# from image_exif import extract_exif
from log import logger
from config import config
import os
import fitz  # PyMuPDF
from pptx import Presentation
from error_decorator import safe_execution


path=config["paths"]["data"]
EMB_PATH = config["paths"]["embeddings"]
META_PATH = config["paths"]["chunks"]
transformer_model=config["model"]["embedding"]
transformer_device=config["device"]["type"]
embedding_batch_size=config["batch"]["embedding_batch_size"]
chunk_size=config["chunking"]["chunk_size"]
chunk_overlap=config["chunking"]["overlap"]
show_progress=config["embedding"]["show_progress"]
retrieve_queries=config["search"]["retrieve_queries"]



def normalize_path(path):
    return os.path.abspath(path)


def sanitize_query(query, max_length=500):
    if not query:
        return None

    # strip whitespace
    query = query.strip()

    # remove excessive whitespace
    query = " ".join(query.split())

    # limit length
    if len(query) > max_length:
        query = query[:max_length]

    # reject empty after cleaning
    if not query:
        return None

    return query


#for images
@safe_execution(component="EXTRACTOR",log_args=True)
def extract_img(path):
    path=normalize_path(path)
    # description_text=str(extract_image(path))
    # exif_text=str(extract_exif(path))
    text=get_image_text_for_embedding(get_file_id_by_path(path))
    mark_processed(path)
    print(text)
    logger.info("Image has been extracted. Path:"+path)
    return text



#for docx
@safe_execution(component="EXTRACTOR",log_args=True)
def extract_docx(path):
    text = docx2txt.process(path)
    mark_processed(normalize_path(path))
    logger.info("Docx has been extracted. Path:"+path)
    return (" ".join(text.split()))

#for pdfs
@safe_execution(component="EXTRACTOR",log_args=True)
def extract_pdf(path):
    doc = fitz.open(path)
    text = ""
    for page in doc:
        text += page.get_text()
    mark_processed(normalize_path(path))
    logger.info("PDF has been extracted. Path:"+path)
    return text

#for pptz
@safe_execution(component="EXTRACTOR",log_args=True)
def extract_pptx(path):
    prs = Presentation(path)
    text = []
    for slide_num, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    content = para.text.strip()
                    if content:
                        text.append(f"[Slide {slide_num+1}] {content}")
    mark_processed(normalize_path(path))
    logger.info("PPTX has been extracted. Path:"+path)
    return "\n".join(text)




#word based chonking
@safe_execution(component="CHUNKER")
def chunk_text(text, chunk_size, overlap):
    words = text.split()
    chunks = []
    step = chunk_size - overlap
    for i in range(0, len(words), step):
        chunk = " ".join(words[i:i + chunk_size])
        if len(chunk.strip()) > 0:
            chunks.append(chunk)
    return chunks




#building the pipeline
@safe_execution(component="PIPELINE", rethrow=True)
def build_pipeline():
    logger.info("Building the pipeline")
    doc_id = 0
    chunker = []

    with os.scandir(path) as es:
        for e in es:
            if not e.is_file():
                continue

            name = e.name.lower()
            if not (name.endswith('.pdf') or name.endswith('.txt') or name.endswith('.pptx') or name.endswith('.docx') or name.endswith('.jpg') or name.endswith('.jpeg') or name.endswith('.png') ):
                continue
            print("Processing:", e.name, e.path)
            mark_processing(normalize_path(e.path))
            text = None
            if name.endswith('.pdf'):
                text = extract_pdf(e.path)
            elif name.endswith('.txt'):
                with open(e.path, encoding='utf-8') as f:
                    text = f.read().strip()
                    mark_processed(normalize_path(e.path))
            elif name.endswith('.pptx'):
                text = extract_pptx(e.path)
            elif name.endswith('.docx'):
                text = extract_docx(e.path)
            elif (name.endswith('.jpg') or name.endswith('.jpeg') or name.endswith('.png') ):
                text=extract_img(e.path)
            if not text or not text.strip():
                continue
            doc_id += 1
            chunks = chunk_text(text,chunk_size,chunk_overlap)
            for i, c in enumerate(chunks):
                chunker.append({
                    "doc_id": doc_id,
                    "chunk_id": i,
                    "doc_name": e.name,
                    "text": c
                })

    print("Total docs:", doc_id)
    print("Total chunks:", len(chunker))
    texts = [c["text"] for c in chunker]
    model = SentenceTransformer(transformer_model,device=transformer_device)
    doc_embeddings = model.encode(
        texts,
        batch_size=embedding_batch_size,
        show_progress_bar=show_progress
    )
    doc_embeddings = doc_embeddings / np.linalg.norm(doc_embeddings, axis=1, keepdims=True)

    np.save(EMB_PATH, doc_embeddings)

    with open(META_PATH, "w", encoding="utf-8") as f:
        json.dump(chunker, f, ensure_ascii=False)

    return model, doc_embeddings, chunker



@safe_execution(component="PIPELINE",rethrow=True)
def load_pipeline():
    print("[+] Loading cached data...")
    logger.info("Loading saved embeddings")
    embeddings = np.load(EMB_PATH)

    with open(META_PATH, "r", encoding="utf-8") as f:
        chunker = json.load(f)

    model = SentenceTransformer(transformer_model, device=transformer_device)

    return model, embeddings, chunker




# Search function based on chunks
@safe_execution(component="SEARCH", default_return=[])
def search(query, model, embeddings, chunker, top_k=5):
    q = sanitize_query(query)
    if not q:
        logger.warning("Empty/invalid query rejected")
        return []
    q = model.encode([q],show_progress_bar=False)[0]
    q = q / np.linalg.norm(q)

    scores = embeddings @ q
    top_indices = np.argsort(scores)[-top_k:][::-1]

    results = []
    for i in top_indices:
        results.append({
            "doc": chunker[i]["doc_name"],
            "score": float(scores[i]),
            "text": chunker[i]["text"][:200]
            })
    logger.info("Chunk search successful. Query:"+q)
    return results

#search function based on docs
@safe_execution(component="SEARCH", default_return=[],log_args=True)
def search_docs(query, model, embeddings, chunker, top_k):
    q = sanitize_query(query)
    if not q:
        logger.warning("Empty/invalid query rejected")
        return []
    # Embed query
    q = model.encode([q],show_progress_bar=False)[0]
    q = q / np.linalg.norm(q)

    # Compute similarity
    scores = embeddings @ q

    # Aggregate scores per document
    doc_scores = {}
    doc_best_chunk = {}

    for i, score in enumerate(scores):
        doc_id = chunker[i]["doc_id"]

        # max pooling
        if doc_id not in doc_scores or score > doc_scores[doc_id]:
            doc_scores[doc_id] = score
            doc_best_chunk[doc_id] = i  # store best chunk index

    # Sort documents
    ranked_docs = sorted(doc_scores.items(), key=lambda x: x[1], reverse=True)

    # Prepare results
    results = []
    for doc_id, score in ranked_docs[:top_k]:
        idx = doc_best_chunk[doc_id]

        results.append({
            "doc": chunker[idx]["doc_name"],
            "score": float(score),
            # "snippet": chunker[idx]["text"][:]
        })

    logger.info("File search successful. Query:"+q)
    return results


# vec = model.encode("hello world")
# print(len(vec))   # should be 384
# docs = [
#     "Python is a programming language",
#     "Cats are animals",
#     "Machine learning is part of AI"
# ]
# doc_embeddings = model.encode(docs)
# print(doc_embeddings.shape)
# print(np.linalg.norm(doc_embeddings[0]))
if __name__ == "__main__":
    logger.info("Main file started")
    if os.path.exists(EMB_PATH) and os.path.exists(META_PATH):
        model, embeddings, chunker = load_pipeline()
    else:
        model, embeddings, chunker = build_pipeline()
    while(True):
        query=input("\nEnter your query. Type E to exit\n")
        if(query=="E"):
            break
        results = search_docs(query, model, embeddings, chunker,retrieve_queries)

        print("\nResults:")
        for r in results:
            print(r)
            print()

