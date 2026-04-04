from sentence_transformers import SentenceTransformer
import numpy as np
import json
import docx2txt

##DOCUMENT CHUNKING
import os
import fitz  # PyMuPDF
from pptx import Presentation
path="data"
EMB_PATH = "embeddings.npy"
META_PATH = "chunks.json"


#for docx
def extract_docx(path):
    text = docx2txt.process(path)
    return (" ".join(text.split()))

#for pdfs
def extract_pdf(path):
    doc = fitz.open(path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text




#for pptz
def extract_pptx(path):
    prs = Presentation(path)
    text = []
    for slide_num, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    content = para.text.strip()
                    if content:
                        text.append(f"[Slide {slide_num+1}] {content}")
    return "\n".join(text)




#word based chonking
def chunk_text(text, chunk_size=400, overlap=100):
    words = text.split()
    chunks = []
    step = chunk_size - overlap
    for i in range(0, len(words), step):
        chunk = " ".join(words[i:i + chunk_size])
        if len(chunk.strip()) > 0:
            chunks.append(chunk)
    return chunks




#building the pipeline
def build_pipeline():
    doc_id = 0
    chunker = []

    with os.scandir(path) as es:
        for e in es:
            if not e.is_file():
                continue

            name = e.name.lower()
            if not (name.endswith('.pdf') or name.endswith('.txt') or name.endswith('.pptx') or name.endswith('.docx')):
                continue
            print("Processing:", e.name)
            text = None
            if name.endswith('.pdf'):
                text = extract_pdf(e.path)
            elif name.endswith('.txt'):
                with open(e.path, encoding='utf-8') as f:
                    text = f.read().strip()
            elif name.endswith('.pptx'):
                text = extract_pptx(e.path)
            elif name.endswith('.docx'):
                text = extract_docx(e.path)
            if not text or not text.strip():
                continue
            doc_id += 1
            chunks = chunk_text(text)
            for i, c in enumerate(chunks):
                chunker.append({
                    "doc_id": doc_id,
                    "chunk_id": i,
                    "doc_name": e.name,
                    "text": c
                })

    print("Total docs:", doc_id)
    print("Total chunks:", len(chunker))
    texts = [c["text"] for c in chunker]
    model = SentenceTransformer('BAAI/bge-small-en-v1.5',device="cuda")
    doc_embeddings = model.encode(
        texts,
        batch_size=32,
        show_progress_bar=True
    )
    doc_embeddings = doc_embeddings / np.linalg.norm(doc_embeddings, axis=1, keepdims=True)

    np.save(EMB_PATH, doc_embeddings)

    with open(META_PATH, "w", encoding="utf-8") as f:
        json.dump(chunker, f, ensure_ascii=False)

    return model, doc_embeddings, chunker




def load_pipeline():
    print("Loading cached data...")

    embeddings = np.load(EMB_PATH)

    with open(META_PATH, "r", encoding="utf-8") as f:
        chunker = json.load(f)

    model = SentenceTransformer('all-MiniLM-L6-v2', device="cuda")

    return model, embeddings, chunker




# Search function based on chunks
def search(query, model, embeddings, chunker, top_k=5):
    q = model.encode([query])[0]
    q = q / np.linalg.norm(q)

    scores = embeddings @ q
    top_indices = np.argsort(scores)[-top_k:][::-1]

    results = []
    for i in top_indices:
        results.append({
            "doc": chunker[i]["doc_name"],
            "score": float(scores[i]),
            "text": chunker[i]["text"][:200]
            })

    return results

#search function based on docs
def search_docs(query, model, embeddings, chunker, top_k=5):
    # Embed query
    q = model.encode([query])[0]
    q = q / np.linalg.norm(q)

    # Compute similarity
    scores = embeddings @ q

    # Aggregate scores per document
    doc_scores = {}
    doc_best_chunk = {}

    for i, score in enumerate(scores):
        doc_id = chunker[i]["doc_id"]

        # max pooling
        if doc_id not in doc_scores or score > doc_scores[doc_id]:
            doc_scores[doc_id] = score
            doc_best_chunk[doc_id] = i  # store best chunk index

    # Sort documents
    ranked_docs = sorted(doc_scores.items(), key=lambda x: x[1], reverse=True)

    # Prepare results
    results = []
    for doc_id, score in ranked_docs[:top_k]:
        idx = doc_best_chunk[doc_id]

        results.append({
            "doc": chunker[idx]["doc_name"],
            "score": float(score),
            # "snippet": chunker[idx]["text"][:]
        })

    return results


# vec = model.encode("hello world")
# print(len(vec))   # should be 384
# docs = [
#     "Python is a programming language",
#     "Cats are animals",
#     "Machine learning is part of AI"
# ]
# doc_embeddings = model.encode(docs)
# print(doc_embeddings.shape)
# print(np.linalg.norm(doc_embeddings[0]))
if __name__ == "__main__":

    if os.path.exists(EMB_PATH) and os.path.exists(META_PATH):
        model, embeddings, chunker = load_pipeline()
    else:
        model, embeddings, chunker = build_pipeline()

    query=input()
    results = search_docs(query, model, embeddings, chunker)

    print("\nResults:")
    for r in results:
        print(r)
        print()
