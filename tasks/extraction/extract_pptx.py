from core.log import logger
from core.errors import safe_execution
from pptx import Presentation
from db.file_repo import mark_processed
from core.utils import normalize_path

#for pptz
@safe_execution(component="EXTRACTOR",log_args=True)
def extract_pptx(path):
    prs = Presentation(path)
    text = []
    for slide_num, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    content = para.text.strip()
                    if content:
                        text.append(f"[Slide {slide_num+1}] {content}")
    mark_processed(normalize_path(path))
    logger.info("PPTX has been extracted. Path:"+path)
    return "\n".join(text)

